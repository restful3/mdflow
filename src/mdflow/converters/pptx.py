"""pptx -> Markdown via python-pptx.

Per slide: title as `## <title>` (or `## Slide N`), body text-frame
paragraphs as a bullet list (2-space indent per paragraph.level), tables
as Markdown tables, and presenter notes as a `> Notes:` blockquote. Shapes
with no text/table (images, graphics) are dropped. No internal try/except.
"""

from __future__ import annotations

import io

from pptx import Presentation

from mdflow.converters.base import (
    ConversionContext,
    ConversionResult,
    ProgressCallback,
)


class PptxConverter:
    name = "pptx-python-pptx"
    formats: tuple[str, ...] = ("pptx",)
    requires_gpu = False

    def can_handle(self, ctx: ConversionContext) -> bool:
        return ctx.format in self.formats

    def convert(self, ctx: ConversionContext, progress: ProgressCallback) -> ConversionResult:
        progress("parse", 10)
        prs = Presentation(io.BytesIO(ctx.data))
        slides = list(prs.slides)
        total = max(len(slides), 1)
        blocks: list[str] = []
        for i, slide in enumerate(slides, start=1):
            blocks.append(_render_slide(slide, i))
            progress("render", 10 + int(80 * i / total))
        markdown = "\n\n".join(blocks).strip()
        progress("done", 100)
        return ConversionResult(markdown=markdown, metadata={"slides": len(slides)})


def _render_slide(slide, index: int) -> str:
    title_shape = slide.shapes.title
    title = title_shape.text.strip() if title_shape is not None else ""
    title_elem = title_shape.element if title_shape is not None else None
    parts: list[str] = [f"## {title}" if title else f"## Slide {index}"]

    for shape in slide.shapes:
        # slide.shapes.title returns a new wrapper each call; compare lxml
        # elements to reliably skip the title shape regardless of wrapper identity.
        if title_elem is not None and shape.element is title_elem:
            continue
        if shape.has_table:
            parts.append(_table_to_md(shape.table))
        elif shape.has_text_frame:
            bullets = _bullets(shape.text_frame)
            if bullets:
                parts.append(bullets)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            quoted = "\n".join(f"> {line}" for line in notes.splitlines())
            parts.append(f"> Notes:\n{quoted}")

    return "\n\n".join(parts)


def _bullets(text_frame) -> str:
    lines: list[str] = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        indent = "  " * (para.level or 0)
        lines.append(f"{indent}- {text}")
    return "\n".join(lines)


def _table_to_md(table) -> str:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header, *body = rows
    out = ["| " + " | ".join(header) + " |"]
    out.append("| " + " | ".join("---" for _ in header) + " |")
    for row in body:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)
