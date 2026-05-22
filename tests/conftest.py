"""Shared pytest fixtures for mdflow tests."""

import io
import shutil
import subprocess
import tempfile
from pathlib import Path

import pytest


@pytest.fixture
def fixtures_dir() -> Path:
    return Path(__file__).parent / "fixtures"


@pytest.fixture
def tmp_cache_dir(tmp_path: Path) -> Path:
    cache = tmp_path / "mdflow_cache"
    cache.mkdir()
    return cache


@pytest.fixture
def sample_docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Document Title", level=1)
    doc.add_heading("Section", level=2)
    p = doc.add_paragraph("Normal text with ")
    p.add_run("bold").bold = True
    p.add_run(" tail.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "First Slide"
    body = s1.placeholders[1].text_frame
    body.text = "Bullet one"
    sub = body.add_paragraph()
    sub.text = "Sub bullet"
    sub.level = 1
    two = body.add_paragraph()
    two.text = "Bullet two"
    s1.notes_slide.notes_text_frame.text = "Presenter note here"

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Second Slide"
    s2.placeholders[1].text_frame.text = "Only bullet"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["name", "score"])
    ws1.append(["alice", 1])
    ws1.append(["bob", 2])
    ws2 = wb.create_sheet("Second")
    ws2.append(["x", "y"])
    ws2.append([10, 20])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_html() -> str:
    return (
        "<!doctype html><html><head><title>T</title></head>\n"
        "<body>\n"
        "<nav>Home About Contact</nav>\n"
        "<article>\n"
        "<h1>Main Heading</h1>\n"
        "<p>First paragraph of the article body with enough text to be detected "
        "as the main content by trafilatura, which needs a reasonable amount of "
        "words before it will treat a block as the article body.</p>\n"
        "<h2>Subsection</h2>\n"
        "<ul><li>Item one</li><li>Item two</li></ul>\n"
        "<table><tr><th>Col A</th><th>Col B</th></tr>"
        "<tr><td>1</td><td>2</td></tr></table>\n"
        "</article>\n"
        "<footer>Copyright 2026</footer>\n"
        "</body></html>"
    )


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    # Larger font sizes let pymupdf4llm's heuristics treat lines as headings.
    page.insert_text((72, 72), "Document Title", fontsize=24)
    page.insert_text((72, 120), "Section One", fontsize=18)
    page.insert_text((72, 150), "First paragraph of body text for the PDF.", fontsize=11)
    page.insert_text((72, 175), "Second paragraph with a bit more content.", fontsize=11)
    data = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def empty_pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    doc.new_page()  # one blank page, no text
    data = doc.tobytes()
    doc.close()
    return data


requires_soffice = pytest.mark.skipif(
    shutil.which("soffice") is None,
    reason="LibreOffice (soffice) not installed",
)


def _soffice_to(src_bytes: bytes, src_ext: str, dst_ext: str) -> bytes:
    """Convert src_bytes (a src_ext document) to dst_ext via soffice headless.

    Used to build legacy binary doc/ppt fixtures from code-generated
    docx/pptx, since there is no pure-Python writer for the legacy
    formats. Caller must be guarded by @requires_soffice.
    """
    soffice = shutil.which("soffice")
    if soffice is None:
        pytest.skip("LibreOffice (soffice) not installed")
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / f"in.{src_ext}"
        src.write_bytes(src_bytes)
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                dst_ext,
                "--outdir",
                tmp,
                f"-env:UserInstallation=file://{Path(tmp) / 'profile'}",
                str(src),
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        return (Path(tmp) / f"in.{dst_ext}").read_bytes()


@pytest.fixture(scope="session")
def sample_doc_bytes() -> bytes:
    from docx import Document

    d = Document()
    d.add_heading("Document Title", level=1)
    d.add_heading("Section One", level=2)
    d.add_paragraph("First paragraph of body text for the doc.")
    buf = io.BytesIO()
    d.save(buf)
    return _soffice_to(buf.getvalue(), "docx", "doc")


@pytest.fixture(scope="session")
def sample_ppt_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = "First Slide"
    s.placeholders[1].text_frame.text = "Bullet one"
    buf = io.BytesIO()
    prs.save(buf)
    return _soffice_to(buf.getvalue(), "pptx", "ppt")
